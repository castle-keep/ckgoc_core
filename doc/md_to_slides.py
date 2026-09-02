#!/usr/bin/env python3
"""
Convert workshop_presentation.md to a PowerPoint (.pptx) that Keynote can open.
Usage: python3 md_to_slides.py workshop_presentation.md -o output.pptx
"""

import sys
import re
from pathlib import Path
from argparse import ArgumentParser
from typing import List
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError as e:
    print(f"Error: missing package — {e}")
    print("Install with: pip install python-pptx lxml")
    sys.exit(1)

# ── Colour palette ──────────────────────────────────────────────────────────
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)   # slide background
C_ACCENT   = RGBColor(0x19, 0x37, 0x73)   # title / heading (dark navy)
C_BODY     = RGBColor(0x1E, 0x1E, 0x1E)   # body text
C_SUB      = RGBColor(0x55, 0x55, 0x55)   # subtext / notes
C_CODE_BG  = RGBColor(0xF3, 0xF4, 0xF6)   # code block bg
C_TH_BG    = RGBColor(0x19, 0x37, 0x73)   # table header bg
C_TH_FG    = RGBColor(0xFF, 0xFF, 0xFF)   # table header fg
C_TR_BG    = RGBColor(0xF0, 0xF4, 0xFF)   # table alt row bg
C_BORDER   = RGBColor(0xCC, 0xCC, 0xCC)   # table border

# ── Slide dimensions (widescreen 13.33 × 7.5 in) ─────────────────────────
W = Inches(13.33)
H = Inches(7.5)
MARGIN = Inches(0.55)


# ─────────────────────────────────────────────────────────────────────────────
# MARKDOWN PARSING
# ─────────────────────────────────────────────────────────────────────────────

def parse_slides(md_path: str) -> List[dict]:
    """Split markdown into list of {title, content, notes} dicts."""
    text = Path(md_path).read_text(encoding='utf-8')

    # Split on `---` horizontal rules (slide separators)
    raw_blocks = re.split(r'\n---\n', text)

    slides = []
    for block in raw_blocks:
        block = block.strip()
        if not block:
            continue

        # Must start with a `# Slide N: ...` heading to be a slide
        m = re.match(r'^# Slide \d+:\s*(.*)', block, re.MULTILINE)
        if not m:
            continue

        # Title: strip the "Slide N:" prefix, keep only the descriptive part
        title = m.group(1).strip()

        # Extract ## Speaker Notes body first (so we can subtract it from body)
        notes_match = re.search(
            r'^## Speaker Notes\n(.*?)(?=\n## |\Z)',
            block, re.MULTILINE | re.DOTALL
        )
        notes = notes_match.group(1).strip() if notes_match else ""

        # Try to find an explicit "## Slide" section
        body_match = re.search(
            r'^## Slide\n(.*?)(?=\n## Speaker Notes\b|\Z)',
            block, re.MULTILINE | re.DOTALL
        )
        if body_match:
            body = body_match.group(1).strip()
        else:
            # No "## Slide" header — grab everything after the title line
            # up to "## Speaker Notes" (or end of block)
            after_title = re.sub(r'^# Slide \d+:.*\n', '', block, count=1)
            if notes_match:
                # Strip the speaker notes section off the end
                notes_start = after_title.find('## Speaker Notes')
                if notes_start != -1:
                    after_title = after_title[:notes_start]
            body = after_title.strip()

        slides.append({'title': title, 'content': body, 'notes': notes})

    return slides


# ─────────────────────────────────────────────────────────────────────────────
# CONTENT BLOCK PARSER
# Splits raw body text into typed blocks for the renderer.
# ─────────────────────────────────────────────────────────────────────────────

def parse_blocks(text: str) -> List[dict]:
    """
    Parse body text into a list of typed blocks:
      {'type': 'bullet',  'level': int, 'text': str}
      {'type': 'table',   'headers': [...], 'rows': [[...]]}
      {'type': 'code',    'text': str}
      {'type': 'heading', 'level': int, 'text': str}
      {'type': 'plain',   'text': str}
    """
    blocks = []
    lines = text.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i]

        # ── fenced code block ──────────────────────────────────────────────
        if line.strip().startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            blocks.append({'type': 'code', 'text': '\n'.join(code_lines)})
            i += 1
            continue

        # ── markdown table ─────────────────────────────────────────────────
        if re.match(r'^\s*\|', line):
            table_lines = []
            while i < len(lines) and re.match(r'^\s*\|', lines[i]):
                table_lines.append(lines[i])
                i += 1
            blocks.append(_parse_table(table_lines))
            continue

        # ── section heading (#### or ###) inside slide body ────────────────
        hm = re.match(r'^(#{2,6})\s+(.*)', line)
        if hm:
            blocks.append({
                'type': 'heading',
                'level': len(hm.group(1)),
                'text': strip_md_inline(hm.group(2).strip())
            })
            i += 1
            continue

        # ── bullet (- item or indented   - item) ──────────────────────────
        bm = re.match(r'^( *)[-*]\s+(.*)', line)
        if bm:
            indent = len(bm.group(1))
            level = 1 if indent < 2 else 2
            blocks.append({
                'type': 'bullet',
                'level': level,
                'text': bm.group(2).strip()
            })
            i += 1
            continue

        # ── blank line ─────────────────────────────────────────────────────
        if not line.strip():
            i += 1
            continue

        # ── plain text ─────────────────────────────────────────────────────
        blocks.append({'type': 'plain', 'text': strip_md_inline(line.strip())})
        i += 1

    return blocks


def _parse_table(table_lines: List[str]) -> dict:
    """Parse markdown table lines into {type, headers, rows}."""
    rows = []
    headers = []
    for raw in table_lines:
        cells = [c.strip() for c in raw.strip().strip('|').split('|')]
        # Skip separator rows (e.g., |---|---|)
        if all(re.match(r'^[-: ]+$', c) for c in cells if c):
            continue
        if not headers:
            headers = [strip_md_inline(c) for c in cells]
        else:
            rows.append([strip_md_inline(c) for c in cells])
    return {'type': 'table', 'headers': headers, 'rows': rows}


def strip_md_inline(text: str) -> str:
    """Remove inline markdown (bold, italic, code backticks) for plain text."""
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*',     r'\1', text)
    text = re.sub(r'`(.+?)`',       r'\1', text)
    return text


# ─────────────────────────────────────────────────────────────────────────────
# RICH TEXT HELPER
# Adds a paragraph with inline bold/code handling.
# ─────────────────────────────────────────────────────────────────────────────

def add_rich_paragraph(tf, text: str, font_size: int, color: RGBColor,
                        bold: bool = False, level: int = 0,
                        space_before: int = 3, space_after: int = 3,
                        monospace: bool = False, first: bool = False):
    """Append a richly-formatted paragraph to a text frame."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)

    # Bullet character for level-1 bullets
    if level == 1:
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•')
    elif level == 2:
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '◦')

    # Split text on **bold** and `code` segments
    segments = re.split(r'(\*\*.*?\*\*|`.*?`)', text)
    for seg in segments:
        run = p.add_run()
        if seg.startswith('**') and seg.endswith('**'):
            run.text = seg[2:-2]
            run.font.bold = True
        elif seg.startswith('`') and seg.endswith('`'):
            run.text = seg[1:-1]
            run.font.bold = False
            run.font.name = 'Courier New'
            run.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
        else:
            run.text = seg
            if bold:
                run.font.bold = True
        run.font.size = Pt(font_size)
        if not (seg.startswith('`') and seg.endswith('`')):
            run.font.color.rgb = color
        if monospace:
            run.font.name = 'Courier New'

    return p


# ─────────────────────────────────────────────────────────────────────────────
# TABLE RENDERER
# ─────────────────────────────────────────────────────────────────────────────

def add_table_shape(slide, headers: List[str], rows: List[List[str]],
                    left, top, width, max_height):
    """Add a formatted pptx table to the slide."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    # Row height: distribute evenly, capped at max_height
    row_h = min(Inches(0.33), max_height // n_rows)
    total_h = row_h * n_rows

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, total_h).table

    # Column widths: equal split
    col_w = width // n_cols
    for ci in range(n_cols):
        tbl.columns[ci].width = col_w

    def _cell_fill(cell, bg: RGBColor):
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = bg

    def _cell_text(cell, text: str, size: int, color: RGBColor,
                   bold: bool = False, align=PP_ALIGN.LEFT):
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        # Handle inline bold/code
        segments = re.split(r'(\*\*.*?\*\*|`.*?`)', text)
        for seg in segments:
            run = p.add_run()
            if seg.startswith('**') and seg.endswith('**'):
                run.text = seg[2:-2]
                run.font.bold = True
            elif seg.startswith('`') and seg.endswith('`'):
                run.text = seg[1:-1]
                run.font.name = 'Courier New'
                run.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
            else:
                run.text = seg
                if bold:
                    run.font.bold = True
            run.font.size = Pt(size)
            run.font.color.rgb = color

    # Header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        _cell_fill(cell, C_TH_BG)
        _cell_text(cell, h, 11, C_TH_FG, bold=True, align=PP_ALIGN.CENTER)

    # Data rows
    for ri, row in enumerate(rows):
        bg = C_TR_BG if ri % 2 == 1 else C_BG
        # Pad or trim row cells to match column count
        padded = (row + [''] * n_cols)[:n_cols]
        for ci, val in enumerate(padded):
            cell = tbl.cell(ri + 1, ci)
            _cell_fill(cell, bg)
            _cell_text(cell, val, 10, C_BODY)

    return tbl


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDER
# ─────────────────────────────────────────────────────────────────────────────

def build_slide(prs: Presentation, slide_data: dict):
    blank = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank)

    # White background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG

    # Accent bar on left edge
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.12), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    # ── Title ─────────────────────────────────────────────────────────────
    title_box = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.28), W - Inches(0.7), Inches(0.82)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide_data['title']
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = C_ACCENT

    # Divider under title
    div = slide.shapes.add_shape(1, Inches(0.35), Inches(1.18), W - Inches(0.7), Pt(2))
    div.fill.solid()
    div.fill.fore_color.rgb = C_ACCENT
    div.line.fill.background()

    # ── Content area ──────────────────────────────────────────────────────
    content_top  = Inches(1.3)
    content_left = Inches(0.5)
    content_w    = W - Inches(0.8)
    content_h    = H - content_top - Inches(0.25)

    blocks = parse_blocks(slide_data['content'])
    if not blocks:
        _write_notes(slide, slide_data['notes'])
        return

    # Check if slide has a table — lay out differently
    has_table  = any(b['type'] == 'table'  for b in blocks)
    has_code   = any(b['type'] == 'code'   for b in blocks)
    text_only  = not has_table and not has_code

    if text_only:
        # Single text box for all bullets/headings/plain
        tb = slide.shapes.add_textbox(content_left, content_top, content_w, content_h)
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for blk in blocks:
            _render_text_block(tf, blk, first)
            first = False
    else:
        # Mixed content: iterate blocks, render each in sequence
        y = content_top
        for blk in blocks:
            remaining_h = H - y - Inches(0.25)
            if blk['type'] == 'table':
                tbl_h = min(remaining_h, Inches(0.35) * (len(blk['rows']) + 1))
                add_table_shape(slide, blk['headers'], blk['rows'],
                                content_left, y, content_w, tbl_h)
                y += tbl_h + Inches(0.15)

            elif blk['type'] == 'code':
                code_lines = blk['text'].split('\n')
                code_h = min(remaining_h, Inches(0.22) * max(len(code_lines), 1) + Inches(0.2))
                code_box = slide.shapes.add_textbox(content_left, y, content_w, code_h)
                code_box.fill.solid()
                code_box.fill.fore_color.rgb = C_CODE_BG
                code_box.line.color.rgb = C_BORDER
                ctf = code_box.text_frame
                ctf.word_wrap = False
                first = True
                for ln in code_lines:
                    p = ctf.paragraphs[0] if first else ctf.add_paragraph()
                    run = p.add_run()
                    run.text = ln
                    run.font.size = Pt(10)
                    run.font.name = 'Courier New'
                    run.font.color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
                    first = False
                y += code_h + Inches(0.1)

            else:
                # Text block (bullet / heading / plain)
                tb = slide.shapes.add_textbox(content_left, y, content_w, Inches(0.5))
                tf = tb.text_frame
                tf.word_wrap = True
                _render_text_block(tf, blk, first=True)
                # Approximate height: 1 line
                line_h = Inches(0.35) if blk['type'] == 'heading' else Inches(0.3)
                y += line_h

    _write_notes(slide, slide_data['notes'])


def _render_text_block(tf, blk: dict, first: bool = False):
    """Render a single text block into an existing text frame."""
    if blk['type'] == 'bullet':
        add_rich_paragraph(
            tf, blk['text'],
            font_size=18 if blk['level'] == 1 else 15,
            color=C_BODY,
            level=blk['level'],
            space_before=5 if blk['level'] == 1 else 2,
            space_after=3,
            first=first,
        )
    elif blk['type'] == 'heading':
        add_rich_paragraph(
            tf, blk['text'],
            font_size=20, color=C_ACCENT, bold=True,
            space_before=12, space_after=4,
            first=first,
        )
    elif blk['type'] == 'plain':
        add_rich_paragraph(
            tf, blk['text'],
            font_size=16, color=C_BODY,
            space_before=4, space_after=2,
            first=first,
        )


def _write_notes(slide, notes_text: str):
    if notes_text:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = notes_text


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    parser = ArgumentParser(description='Convert workshop_presentation.md → .pptx (Keynote-compatible)')
    parser.add_argument('input',  help='Input .md file')
    parser.add_argument('-o', '--output', default='presentation.pptx', help='Output .pptx file')
    args = parser.parse_args()

    if not Path(args.input).exists():
        print(f"Error: {args.input} not found")
        sys.exit(1)

    print(f"Parsing {args.input}…")
    slides = parse_slides(args.input)
    print(f"Found {len(slides)} slides")

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for idx, slide_data in enumerate(slides, 1):
        build_slide(prs, slide_data)
        print(f"  [{idx:02d}/{len(slides)}] {slide_data['title']}")

    prs.save(args.output)
    print(f"\n✓  Saved → {args.output}")
    print("   Open in Keynote or PowerPoint.")


if __name__ == '__main__':
    main()
